"""
M.A.S.H Presentation Generator
Run with: /home/bluni/miniconda3/bin/python make_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────
NAVY      = RGBColor(0x0A, 0x1F, 0x44)   # deep navy – slide backgrounds
TEAL      = RGBColor(0x00, 0x96, 0x99)   # accent teal
LIGHT_TEAL= RGBColor(0xC8, 0xF4, 0xF5)  # very light teal for text boxes
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
PALE_GREY = RGBColor(0xF2, 0xF6, 0xFA)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
MID_GREY  = RGBColor(0x66, 0x66, 0x66)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
ORANGE    = RGBColor(0xE6, 0x7E, 0x22)
RED       = RGBColor(0xC0, 0x39, 0x2B)
AMBER     = RGBColor(0xF3, 0x9C, 0x12)

# Slide dimensions (16:9 widescreen)
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank


# ── Helper functions ─────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def text_box(slide, text, x, y, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False,
             v_anchor=None):
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    if v_anchor:
        txBox.text_frame.word_wrap = wrap
        txBox.text_frame._txBody.get_or_add_bodyPr().set('anchor', v_anchor)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=14, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def navy_header(slide, title, subtitle=None):
    """Full-width navy header bar at top of slide."""
    rect(slide, 0, 0, 13.33, 1.1, fill=NAVY)
    text_box(slide, title, 0.3, 0.05, 12, 0.65,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        text_box(slide, subtitle, 0.3, 0.65, 12, 0.38,
                 font_size=14, color=LIGHT_TEAL, align=PP_ALIGN.LEFT)
    # teal accent line under header
    rect(slide, 0, 1.1, 13.33, 0.04, fill=TEAL)


def card(slide, x, y, w, h, bg=PALE_GREY, border=TEAL):
    return rect(slide, x, y, w, h, fill=bg, line=border, line_w=Pt(1.5))


def bullet_block(slide, items, x, y, w, h,
                 bg=PALE_GREY, border=TEAL,
                 title=None, title_color=NAVY,
                 body_color=DARK_GREY, font_size=13):
    card(slide, x, y, w, h, bg=bg, border=border)
    ty = y + 0.1
    if title:
        text_box(slide, title, x + 0.15, ty, w - 0.3, 0.38,
                 font_size=15, bold=True, color=title_color)
        ty += 0.38
    txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(ty),
                                     Inches(w - 0.3), Inches(h - (ty - y) - 0.1))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = body_color
    return txBox


def teal_badge(slide, label, x, y, w=1.6, h=0.35):
    rect(slide, x, y, w, h, fill=TEAL)
    text_box(slide, label, x, y, w, h,
             font_size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)


def flow_arrow(slide, x1, y, x2):
    """Horizontal arrow from x1 to x2 at height y (inches)."""
    from pptx.util import Inches, Pt
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(x1), Inches(y), Inches(x2), Inches(y)
    )
    connector.line.color.rgb = TEAL
    connector.line.width = Pt(2)
    return connector


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════
sl = add_slide()

# Full background navy
rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)

# Teal diagonal accent stripe (wide rectangle rotated via shape)
accent = rect(sl, -1, 5.2, 15, 0.5, fill=TEAL)

# White background card in centre
rect(sl, 1.2, 1.5, 10.9, 4.4, fill=WHITE)

# Title
text_box(sl, "M.A.S.H", 1.6, 1.7, 10, 1.4,
         font_size=72, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Teal underline
rect(sl, 4.0, 3.1, 5.3, 0.07, fill=TEAL)

# Subtitle
text_box(sl, "Medical Assistant & Services Hub", 1.6, 3.22, 10, 0.6,
         font_size=24, bold=False, color=TEAL, align=PP_ALIGN.CENTER)

# Tagline
text_box(sl,
         "6 Specialized AI Agents · Multi-Room Secure Network · Zero Admin Friction",
         1.6, 3.85, 10, 0.5,
         font_size=15, color=MID_GREY, align=PP_ALIGN.CENTER)

# Bottom label
text_box(sl, "Band of Agents SDK  ·  Supabase  ·  LangGraph  ·  Gemini",
         0, 6.3, 13.33, 0.5,
         font_size=13, color=LIGHT_TEAL, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
navy_header(sl, "The Problem", "Healthcare is drowning in administrative overhead")

problems = [
    ("Manual Triage",       "Schedulers hand-match symptoms to specialties.\nSlow, inconsistent, prone to bias.",             RED),
    ("EHR Click-Fatigue",   "Doctors spend 15–20 min per patient clicking\nthrough fragmented history screens.",              ORANGE),
    ("Pharmacy Blindspots", "Prescriptions written without checking stock.\nPatients turned away at the pharmacy.",           AMBER),
    ("Navigation Anxiety",  "Patients wander hallways asking staff for\ndirections, creating lobby congestion.",              TEAL),
]

cx = 0.4
for title, desc, color in problems:
    rect(sl, cx, 1.35, 2.9, 5.6, fill=PALE_GREY, line=color, line_w=Pt(3))
    rect(sl, cx, 1.35, 2.9, 0.55, fill=color)
    text_box(sl, title, cx + 0.1, 1.38, 2.7, 0.5,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(sl, desc, cx + 0.15, 2.0, 2.6, 4.7,
             font_size=13, color=DARK_GREY, wrap=True)
    cx += 3.2

# Bottom insight
rect(sl, 0.3, 7.0, 12.73, 0.38, fill=NAVY)
text_box(sl, "Result: Clinicians lose hours. Patients get frustrated. Errors happen.",
         0.3, 7.02, 12.73, 0.35,
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — What is M.A.S.H
# ════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
navy_header(sl, "What is M.A.S.H?", "A decentralized AI orchestration platform for healthcare clinics")

text_box(sl,
         "M.A.S.H coordinates 6 specialized AI agents across a secure, multi-room virtual "
         "network to automate clinical summaries, doctor matching, and prescription stock checks — "
         "eliminating administrative burden with a minimal, jargon-free interface.",
         0.4, 1.3, 12.5, 0.95,
         font_size=15, color=DARK_GREY, wrap=True)

# Three pillars
pillars = [
    ("For Patients",      TEAL,
     ["Remote booking via Mobile Interface",
      "Symptom chatbot assigns right specialist",
      "In-clinic wayfinding to doctor's room",
      "No paper forms, no queuing at reception"]),
    ("For Doctors",       NAVY,
     ["Auto-compiled patient summary on arrival",
      "Medical history, labs & surgeries in one view",
      "Prescription triggers instant stock check",
      "Alert if drug unavailable → suggest alternative"]),
    ("For Pharmacy",      GREEN,
     ["Real-time prescription orders from doctor",
      "Stock deducted automatically on fulfillment",
      "High-demand drugs flagged for pre-ordering",
      "Human-in-loop approval for out-of-stock items"]),
]

px = 0.4
for title, color, items in pillars:
    rect(sl, px, 2.35, 4.0, 4.8, fill=PALE_GREY, line=color, line_w=Pt(2.5))
    rect(sl, px, 2.35, 4.0, 0.5, fill=color)
    text_box(sl, title, px + 0.1, 2.38, 3.8, 0.44,
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    ty = 2.95
    for item in items:
        text_box(sl, "▸  " + item, px + 0.2, ty, 3.6, 0.45,
                 font_size=13, color=DARK_GREY)
        ty += 0.52
    px += 4.3


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — System Architecture Overview
# ════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
navy_header(sl, "System Architecture", "Multi-room, event-driven agent mesh on the Band of Agents SDK")

# Legend
text_box(sl, "Interfaces", 0.4, 1.25, 2.0, 0.3, font_size=11, bold=True, color=NAVY)
text_box(sl, "Virtual Rooms (Event Buses)", 3.0, 1.25, 4.0, 0.3, font_size=11, bold=True, color=TEAL)
text_box(sl, "AI Agents", 8.5, 1.25, 2.5, 0.3, font_size=11, bold=True, color=GREEN)
text_box(sl, "Database", 11.3, 1.25, 2.0, 0.3, font_size=11, bold=True, color=ORANGE)

# ─── Interfaces column ───
rect(sl, 0.2, 1.65, 2.3, 0.75, fill=NAVY)
text_box(sl, "Mobile App\n(Patient)", 0.2, 1.65, 2.3, 0.75,
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(sl, 0.2, 2.55, 2.3, 0.75, fill=NAVY)
text_box(sl, "Doctor\nDashboard", 0.2, 2.55, 2.3, 0.75,
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(sl, 0.2, 3.45, 2.3, 0.75, fill=NAVY)
text_box(sl, "Pharmacist\nDashboard", 0.2, 3.45, 2.3, 0.75,
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(sl, 0.2, 4.35, 2.3, 0.75, fill=NAVY)
text_box(sl, "Express.js\nAPI Server", 0.2, 4.35, 2.3, 0.75,
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─── Rooms column ───
rooms = [
    "Patient-Management-Room",
    "Reception-Navigation-Room",
    "Clinical-Consult-Room",
    "Doctor-Dashboard-Room",
    "Pharmacy-Inventory-Room",
    "Telemetry-Audit-Room",
]
ry = 1.55
for room in rooms:
    rect(sl, 3.0, ry, 3.8, 0.58, fill=RGBColor(0xE8,0xF8,0xFA), line=TEAL, line_w=Pt(1.5))
    text_box(sl, room, 3.05, ry + 0.08, 3.7, 0.44, font_size=11, color=TEAL, bold=True)
    ry += 0.65

# ─── Agents column ───
agents = [
    ("Registration Agent",      GREEN),
    ("Patient Mgmt Agent",      GREEN),
    ("Navigation Agent",        GREEN),
    ("Summary Agent",           GREEN),
    ("Medicine Mgmt Agent",     GREEN),
    ("Stock Mgmt Agent",        GREEN),
]
ay = 1.55
for name, color in agents:
    rect(sl, 7.3, ay, 3.3, 0.56, fill=RGBColor(0xE8,0xF9,0xEE), line=color, line_w=Pt(1.5))
    text_box(sl, name, 7.35, ay + 0.08, 3.2, 0.42, font_size=12, color=color, bold=True)
    ay += 0.65

# ─── Database column ───
rect(sl, 11.0, 1.65, 2.1, 1.1, fill=RGBColor(0xFF,0xF0,0xE0), line=ORANGE, line_w=Pt(1.5))
text_box(sl, "Supabase\n(PostgreSQL\n+ RLS)", 11.0, 1.7, 2.1, 1.1,
         font_size=12, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

rect(sl, 11.0, 2.95, 2.1, 0.75, fill=PALE_GREY, line=MID_GREY, line_w=Pt(1))
text_box(sl, "Band.ai\nSDK", 11.0, 2.95, 2.1, 0.75,
         font_size=12, color=DARK_GREY, bold=True, align=PP_ALIGN.CENTER)

rect(sl, 11.0, 3.85, 2.1, 0.75, fill=PALE_GREY, line=MID_GREY, line_w=Pt(1))
text_box(sl, "Gemini\nLLM", 11.0, 3.85, 2.1, 0.75,
         font_size=12, color=DARK_GREY, bold=True, align=PP_ALIGN.CENTER)

# Arrow labels
text_box(sl, "←  Events  →", 2.6, 3.8, 0.5, 4.0,
         font_size=9, color=MID_GREY, italic=True)
text_box(sl, "←  Events  →", 6.7, 3.8, 0.5, 4.0,
         font_size=9, color=MID_GREY, italic=True)
text_box(sl, "← DB / SDK →", 10.5, 3.8, 0.5, 4.0,
         font_size=9, color=MID_GREY, italic=True)

# Bottom note
rect(sl, 0.2, 6.8, 12.9, 0.48, fill=NAVY)
text_box(sl,
         "Rooms are persistent event buses. Agents subscribe to rooms & communicate exclusively via typed events.",
         0.3, 6.82, 12.7, 0.44,
         font_size=12, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — The 6 Agents
# ════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
navy_header(sl, "The 6 AI Agents", "Single-responsibility agents, each owning a clinical domain")

agents_data = [
    ("1  Registration Agent",
     "Smart Clinical Matchmaker",
     ["Maintains real-time doctor directory & schedules",
      "Maps symptoms → specialty (NLP rules + LLM)",
      "Assigns best available specialist automatically"],
     TEAL,
     "SUBSCRIBES: REQUEST_DOCTOR_MATCH\nPUBLISHES: DOCTOR_ASSIGNED"),
    ("2  Patient Mgmt Agent",
     "Queue & Booking Orchestrator",
     ["Processes remote appointments from mobile app",
      "Manages live queues with check-ins & rescheduling",
      "Front-line LLM chat interface for the patient"],
     NAVY,
     "SUBSCRIBES: BOOK_APPOINTMENT, PATIENT_CHECK_IN\nPUBLISHES: QUEUE_UPDATED, RESCHEDULE_CONFIRMED"),
    ("3  Navigation Agent",
     "Wayfinding Assistant",
     ["Monitors doctor-to-room allocations in real time",
      "Calculates step-by-step directions on check-in",
      "Handles room relocations dynamically"],
     GREEN,
     "SUBSCRIBES: NAVIGATE_TO_ROOM, DOCTOR_ROOM_CHANGE\nPUBLISHES: NAVIGATION_DIRECTIONS"),
    ("4  Summary Agent",
     "Clinical Record Aggregator",
     ["Queries allergies, vitals, lab results, surgeries",
      "LLM compiles a structured patient brief",
      "Saves summary back to Supabase"],
     ORANGE,
     "SUBSCRIBES: SUMMARIZE_PATIENT_HISTORY\nPUBLISHES: PATIENT_HISTORY_COMPILED"),
    ("5  Medicine Mgmt Agent",
     "Prescription Safety Auditor",
     ["Intercepts every written prescription",
      "Cross-references live inventory stock",
      "Human-in-loop for out-of-stock alternatives"],
     RED,
     "SUBSCRIBES: PRESCRIPTION_WRITTEN\nPUBLISHES: PREPARE_MEDICINE, ALTERNATIVE_REQUESTED"),
    ("6  Stock Mgmt Agent",
     "Inventory & Logistics Controller",
     ["Tracks real-time medicine stock levels",
      "Deducts stock on each fulfilled prescription",
      "Flags high-demand drugs for proactive reorder"],
     AMBER,
     "SUBSCRIBES: ROUTE_TO_PHARMA\nPUBLISHES: STOCK_DEMAND_ALERT, TRIGGER_REORDER"),
]

col = 0
row = 0
for i, (title, role, bullets, color, events) in enumerate(agents_data):
    x = 0.3 + col * 4.3
    y = 1.3 + row * 3.0
    rect(sl, x, y, 4.05, 2.85, fill=PALE_GREY, line=color, line_w=Pt(2))
    rect(sl, x, y, 4.05, 0.48, fill=color)
    text_box(sl, title, x + 0.1, y + 0.04, 3.85, 0.42,
             font_size=13, bold=True, color=WHITE)
    text_box(sl, role, x + 0.1, y + 0.52, 3.85, 0.32,
             font_size=11, bold=True, color=color, italic=True)
    by = y + 0.88
    for b in bullets:
        text_box(sl, "• " + b, x + 0.1, by, 3.85, 0.38, font_size=10.5, color=DARK_GREY)
        by += 0.4
    text_box(sl, events, x + 0.1, y + 2.45, 3.85, 0.36,
             font_size=8.5, color=MID_GREY, italic=True)
    col += 1
    if col == 3:
        col = 0
        row += 1


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Virtual Room Architecture
# ════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
navy_header(sl, "Virtual Room Architecture", "Persistent secure event buses — one concern per room")

rooms_info = [
    ("Patient-Management-Room",
     "PatientManagementAgent\nRegistrationAgent",
     "Booking, rescheduling, doctor availability, appointment queue",
     TEAL),
    ("Reception-Navigation-Room",
     "PatientNavigationAgent\nRegistrationAgent",
     "Doctor matching at front desk, wayfinding, room-change updates",
     NAVY),
    ("Clinical-Consult-Room",
     "MedicineManagementAgent\nSummaryAgent",
     "Patient history compilation, prescription safety handoff",
     GREEN),
    ("Doctor-Dashboard-Room",
     "DoctorAssistantAgent\nSummaryAgent",
     "Clinical summaries, booking notifications to doctor",
     ORANGE),
    ("Pharmacy-Inventory-Room",
     "MedicineManagementAgent\nStockManagementAgent",
     "Stock checks, order routing, reorder alerts",
     RED),
    ("Telemetry-Audit-Room",
     "TelemetryAgent + All Agents",
     "Cross-cutting audit log, observability, compliance trail",
     AMBER),
]

rx = 0.3
ry = 1.3
col = 0
for name, agents_in_room, purpose, color in rooms_info:
    rect(sl, rx, ry, 4.1, 2.7, fill=PALE_GREY, line=color, line_w=Pt(2))
    rect(sl, rx, ry, 4.1, 0.42, fill=color)
    text_box(sl, name, rx + 0.1, ry + 0.04, 3.9, 0.36,
             font_size=12, bold=True, color=WHITE)
    text_box(sl, "Agents:  " + agents_in_room,
             rx + 0.1, ry + 0.48, 3.9, 0.7,
             font_size=10.5, color=NAVY, bold=True)
    text_box(sl, purpose, rx + 0.1, ry + 1.22, 3.9, 1.35,
             font_size=11, color=DARK_GREY, wrap=True)
    col += 1
    rx += 4.35
    if col == 3:
        col = 0
        rx = 0.3
        ry += 3.0

# Design principle note
rect(sl, 0.3, 7.05, 12.73, 0.35, fill=NAVY)
text_box(sl,
         "Rooms are initialised once at startup. Room IDs are stored in .env.rooms — no dynamic creation.",
         0.4, 7.07, 12.5, 0.3,
         font_size=11, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Workflow: Patient Booking Journey
# ════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
navy_header(sl, "Workflow 1 — Symptom Triage & Booking",
            "From patient symptom description to confirmed appointment in seconds")

steps = [
    ("Patient",          "\"Chest pain &\npalpitations\"",              NAVY),
    ("Patient\nMgmt\nAgent",   "LLM parses\nintent &\ntriggers tool",   TEAL),
    ("Patient-Mgmt\nRoom",     "Broadcasts\nREQUEST_\nDOCTOR_MATCH",   TEAL),
    ("Registration\nAgent",    "Maps symptoms\n→ Cardiology\nDB query", GREEN),
    ("Supabase",               "Returns\navailable\nCardiologists",      ORANGE),
    ("Registration\nAgent",    "Broadcasts\nDOCTOR_\nASSIGNED",         GREEN),
    ("Registration\nAgent",    "Broadcasts\nBOOKING_\nCONFIRMED",       GREEN),
    ("Patient",                "\"Dr. Chen,\n09:00 AM\nconfirmed\"",     NAVY),
]

sx = 0.3
sy = 1.55
box_w = 1.45
box_h = 1.55
gap = 0.2

for i, (actor, label, color) in enumerate(steps):
    rect(sl, sx, sy, box_w, box_h, fill=RGBColor(0xE8,0xF8,0xFA) if color == TEAL
         else (RGBColor(0xE8,0xF9,0xEE) if color == GREEN
         else (RGBColor(0xFF,0xF0,0xE0) if color == ORANGE
         else PALE_GREY)),
         line=color, line_w=Pt(2))
    text_box(sl, actor, sx, sy, box_w, 0.42,
             font_size=10, bold=True, color=WHITE if color == NAVY else color,
             align=PP_ALIGN.CENTER)
    rect(sl, sx, sy, box_w, 0.38, fill=color)
    text_box(sl, actor, sx, sy + 0.03, box_w, 0.35,
             font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(sl, label, sx + 0.05, sy + 0.45, box_w - 0.1, box_h - 0.55,
             font_size=10, color=DARK_GREY, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        # Arrow
        ax = sx + box_w + gap / 2
        text_box(sl, "→", ax - 0.05, sy + box_h / 2 - 0.2, 0.2, 0.4,
                 font_size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    sx += box_w + gap + 0.1

# Symptom classification table
rect(sl, 0.3, 3.3, 12.73, 0.35, fill=NAVY)
text_box(sl, "Symptom Classification Rules (Registration Agent — LangGraph triage node)",
         0.4, 3.32, 12.5, 0.3, font_size=12, bold=True, color=WHITE)

rows = [
    ('"chest pain", "heart", "cardio"',   "→  Cardiology"),
    ('"fever", "child", "pediatric"',     "→  Pediatrics"),
    ("All other inputs",                  "→  General Practice"),
]
ry = 3.75
for kw, spec in rows:
    rect(sl, 0.3, ry, 7.0, 0.42, fill=PALE_GREY, line=TEAL, line_w=Pt(1))
    text_box(sl, kw,   0.5,  ry + 0.06, 6.0, 0.32, font_size=12, color=DARK_GREY, italic=True)
    rect(sl, 7.3, ry, 5.73, 0.42, fill=PALE_GREY, line=GREEN, line_w=Pt(1))
    text_box(sl, spec, 7.5,  ry + 0.06, 5.0, 0.32, font_size=12, bold=True, color=GREEN)
    ry += 0.52

# asyncio Future pattern note
rect(sl, 0.3, 5.12, 12.73, 1.05, fill=RGBColor(0xFF,0xF8,0xE8), line=AMBER, line_w=Pt(1.5))
text_box(sl, "Proxy Tool Pattern (Key Innovation)",
         0.5, 5.17, 12.3, 0.32, font_size=12, bold=True, color=AMBER)
text_box(sl,
         "PatientManagementAgent creates an asyncio.Future keyed by requestId → broadcasts event to room → "
         "awaits the Future (non-blocking).\nRegistrationAgent executes the DB query, broadcasts response "
         "with same requestId → PatientManagementAgent resolves Future → LLM resumes.",
         0.5, 5.52, 12.3, 0.6,
         font_size=11.5, color=DARK_GREY, wrap=True)

rect(sl, 0.3, 6.3, 12.73, 0.55, fill=PALE_GREY)
text_box(sl,
         "Result: LLM agent (PatientMgmtAgent) is cleanly decoupled from DB agent (RegistrationAgent). "
         "Timeout = 10 s → graceful failure message to patient, not a crash.",
         0.5, 6.33, 12.3, 0.48,
         font_size=11, color=DARK_GREY, wrap=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Workflow: In-Clinic Arrival & Navigation
# ════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
navy_header(sl, "Workflow 2 — In-Clinic Arrival & Wayfinding",
            "Patient checks in → instant room directions + clinical brief")

# Left: Navigation flow
rect(sl, 0.3, 1.25, 6.0, 5.95, fill=PALE_GREY, line=TEAL, line_w=Pt(1.5))
text_box(sl, "Patient Navigation Flow", 0.4, 1.3, 5.8, 0.4,
         font_size=14, bold=True, color=TEAL)

nav_steps = [
    ("Patient arrives & checks in on Mobile App", NAVY),
    ("PatientManagementAgent broadcasts NAVIGATE_TO_ROOM\n{patientId, doctorId, currentLocation}", TEAL),
    ("PatientNavigationAgent looks up doctor-room index\n(doc-1 → Room 302, 3rd Floor)", GREEN),
    ("Broadcasts NAVIGATION_DIRECTIONS back to Reception-Navigation-Room", GREEN),
    ("Mobile App displays: \"Go to elevator → 3rd Floor → Room 302\"", NAVY),
]
ny = 1.78
for step, color in nav_steps:
    rect(sl, 0.5, ny, 5.6, 0.72, fill=WHITE, line=color, line_w=Pt(1.5))
    text_box(sl, step, 0.65, ny + 0.08, 5.3, 0.6, font_size=11, color=DARK_GREY, wrap=True)
    ny += 0.85
    if ny < 6.9:
        text_box(sl, "↓", 3.1, ny - 0.12, 0.3, 0.3, font_size=16, bold=True, color=TEAL,
                 align=PP_ALIGN.CENTER)

# Right: Clinical Summary flow
rect(sl, 6.8, 1.25, 6.2, 5.95, fill=PALE_GREY, line=ORANGE, line_w=Pt(1.5))
text_box(sl, "Clinical Summary Compilation", 6.9, 1.3, 6.0, 0.4,
         font_size=14, bold=True, color=ORANGE)

sum_steps = [
    ("Doctor triggers SUMMARIZE_PATIENT_HISTORY\n{patientId}", ORANGE),
    ("SummaryAgent queries Supabase medical_records\nby patient UUID", ORANGE),
    ("Extracts: chronic conditions, allergies,\ntest results, surgical history", DARK_GREY),
    ("LLM (Gemini) compiles structured markdown brief", GREEN),
    ("Summary saved to DB + broadcast PATIENT_HISTORY_COMPILED\nto Clinical-Consult-Room & Doctor-Dashboard-Room", GREEN),
]
sy2 = 1.78
for step, color in sum_steps:
    rect(sl, 7.0, sy2, 5.8, 0.72, fill=WHITE, line=ORANGE, line_w=Pt(1.5))
    text_box(sl, step, 7.15, sy2 + 0.08, 5.5, 0.6, font_size=11, color=DARK_GREY, wrap=True)
    sy2 += 0.85
    if sy2 < 6.9:
        text_box(sl, "↓", 9.75, sy2 - 0.12, 0.3, 0.3, font_size=16, bold=True, color=ORANGE,
                 align=PP_ALIGN.CENTER)

# Data table
rect(sl, 0.3, 7.05, 12.73, 0.35, fill=NAVY)
text_box(sl,
         "DB record_type mapping: chronic_condition · allergy (JSON) · test_result (JSON) · surgical_history (JSON)",
         0.4, 7.07, 12.5, 0.3, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Workflow: Prescription & Pharmacy Routing
# ════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
navy_header(sl, "Workflow 3 — Prescription & Pharmacy Routing",
            "Every prescription is safety-checked before the patient leaves the room")

# Main flow columns
col_titles = ["Doctor writes\nprescription", "Medicine Mgmt\nAgent", "Stock Mgmt\nAgent + DB",
              "IN STOCK →\nPharmacist Room", "OUT OF STOCK →\nDoctor Alert"]
col_colors = [NAVY, TEAL, GREEN, GREEN, RED]
cx2 = 0.3
for i, (ct, cc) in enumerate(zip(col_titles, col_colors)):
    rect(sl, cx2, 1.3, 2.4, 0.75, fill=cc)
    text_box(sl, ct, cx2, 1.3, 2.4, 0.75,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    cx2 += 2.6

# In-stock path
rect(sl, 0.3, 2.2, 2.4, 4.95, fill=PALE_GREY, line=NAVY, line_w=Pt(1))
text_box(sl, "PRESCRIPTION_WRITTEN\n{patientId, medicine,\ndosage, duration}",
         0.4, 2.3, 2.2, 1.0, font_size=11, color=DARK_GREY)

rect(sl, 2.9, 2.2, 2.4, 4.95, fill=PALE_GREY, line=TEAL, line_w=Pt(1))
text_box(sl, "Query Supabase\nfor stock level\n\nRoute based\non availability",
         3.0, 2.3, 2.2, 1.8, font_size=11, color=DARK_GREY)

rect(sl, 5.5, 2.2, 2.4, 4.95, fill=PALE_GREY, line=GREEN, line_w=Pt(1))
text_box(sl, "Checks stock table\nReturns level\n\nDeducts stock\nIncrements\nusage counter",
         5.6, 2.3, 2.2, 2.5, font_size=11, color=DARK_GREY)

# In-stock branch
rect(sl, 8.1, 2.2, 2.4, 2.35, fill=RGBColor(0xE8,0xF9,0xEE), line=GREEN, line_w=Pt(1.5))
text_box(sl, "✓ In Stock", 8.1, 2.25, 2.4, 0.35,
         font_size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
text_box(sl, "Broadcast PREPARE_MEDICINE\nto Pharmacist-Dashboard\n\nPharmacist prints\npreparation alert",
         8.15, 2.65, 2.3, 1.75, font_size=10.5, color=DARK_GREY)

# Out-of-stock branch
rect(sl, 8.1, 4.7, 2.4, 2.35, fill=RGBColor(0xFF,0xEB,0xEB), line=RED, line_w=Pt(1.5))
text_box(sl, "✗ Out of Stock", 8.1, 4.75, 2.4, 0.35,
         font_size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)
text_box(sl, "Broadcast ALTERNATIVE_MEDICINE_\nREQUESTED to Doctor\n\nSystem PAUSES\n→ Human-in-loop\napproval",
         8.15, 5.15, 2.3, 1.75, font_size=10.5, color=DARK_GREY)

# High demand alert
rect(sl, 10.7, 2.2, 2.4, 4.95, fill=RGBColor(0xFF,0xF8,0xE8), line=AMBER, line_w=Pt(1.5))
text_box(sl, "High Demand Alert", 10.7, 2.25, 2.4, 0.35,
         font_size=11, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
text_box(sl,
         "If medicine prescribed\n≥ 2 times this shift:\n\nBroadcast\nSTOCK_DEMAND_ALERT\nto Pharmacist\n\n→ Prepare extra\n  stock in advance",
         10.75, 2.65, 2.3, 4.35, font_size=10.5, color=DARK_GREY)

# Arrows
for ax2 in [2.7, 5.3, 7.9]:
    text_box(sl, "→", ax2, 4.4, 0.25, 0.4, font_size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
text_box(sl, "↑ usage >= 2", 10.45, 3.0, 0.3, 0.8, font_size=9, color=AMBER, italic=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — Agent Handoff Sequence (end-to-end patient visit)
# ════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
navy_header(sl, "End-to-End: Full Patient Visit Flow",
            "All 6 agents coordinating across rooms for a single patient visit")

phases = [
    ("REMOTE BOOKING",      NAVY,   [
        "Patient describes symptoms on mobile chatbot",
        "PatientMgmtAgent triggers REQUEST_DOCTOR_MATCH",
        "RegistrationAgent matches Cardiology → Dr. Chen",
        "Appointment written to Supabase → BOOKING_CONFIRMED",
    ]),
    ("IN-CLINIC ARRIVAL",   TEAL,   [
        "Patient checks in via Mobile App",
        "NavigationAgent sends directions: Room 205",
        "SummaryAgent compiles history from Supabase",
        "Summary displayed on Doctor Dashboard",
    ]),
    ("CONSULTATION",        GREEN,  [
        "Doctor reviews AI-compiled patient brief",
        "Prescribes Rare-Antibiotic",
        "MedicineAgent checks stock → OUT OF STOCK",
        "Doctor notified → selects Ibuprofen instead",
    ]),
    ("PHARMACY",            ORANGE, [
        "MedicineAgent confirms Ibuprofen in stock",
        "StockAgent deducts 1 unit from inventory",
        "PREPARE_MEDICINE sent to Pharmacist Dashboard",
        "Patient prescription fulfilled & sent",
    ]),
]

px2 = 0.3
for phase, color, steps2 in phases:
    rect(sl, px2, 1.3, 3.1, 5.9, fill=PALE_GREY, line=color, line_w=Pt(2))
    rect(sl, px2, 1.3, 3.1, 0.48, fill=color)
    text_box(sl, phase, px2 + 0.05, 1.33, 3.0, 0.42,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    sy3 = 1.9
    for s in steps2:
        text_box(sl, "▸  " + s, px2 + 0.15, sy3, 2.85, 0.56,
                 font_size=12, color=DARK_GREY, wrap=True)
        sy3 += 0.65
    px2 += 3.25

# Central insight
rect(sl, 0.3, 7.0, 12.73, 0.4, fill=NAVY)
text_box(sl,
         "Zero manual handoffs. All agents react to events in real-time. "
         "Doctors never touch a scheduling tool. Pharmacy never misses a prescription.",
         0.4, 7.02, 12.5, 0.36,
         font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — Architecture Design Principles
# ════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
navy_header(sl, "Architecture Design Principles", "The engineering decisions that make M.A.S.H safe and scalable")

principles = [
    ("Persistent Rooms,\nDynamic Sessions",
     TEAL,
     "Rooms are long-lived event buses, initialised once at server startup. "
     "Individual workflows are transaction-bound with a UUID (requestId). "
     "Avoids API overhead from dynamic room creation. Agents always know which room to find.",
     ),
    ("Static Graphs,\nDynamic Execution",
     GREEN,
     "All backend agents use LangGraph StateGraphs compiled at startup. "
     "Graph layout (nodes, edges, safety checks) is fixed — cannot be bypassed by LLM hallucination. "
     "Runtime invocation creates isolated execution threads per patient request.",
     ),
    ("Proxy Tool Pattern\n(Async Future Bridge)",
     NAVY,
     "LLM agents call 'proxy tools' that create asyncio.Futures. "
     "The Future is stored by requestId; the LLM awaits non-blocking. "
     "Backend DB agents resolve the Future on response. Timeout = 10 s → graceful failure.",
     ),
    ("Human-in-the-Loop\nfor Clinical Safety",
     RED,
     "Prescription workflows pause execution if a drug is out of stock. "
     "Resumes only after doctor provides explicit alternative. "
     "Structural guarantee — the safety gate cannot be removed without changing the graph.",
     ),
]

px3 = 0.3
for title, color, desc in principles:
    rect(sl, px3, 1.3, 3.1, 5.9, fill=PALE_GREY, line=color, line_w=Pt(2.5))
    rect(sl, px3, 1.3, 3.1, 0.75, fill=color)
    text_box(sl, title, px3 + 0.08, 1.35, 2.94, 0.65,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(sl, desc, px3 + 0.15, 2.15, 2.82, 4.85,
             font_size=12, color=DARK_GREY, wrap=True)
    px3 += 3.25


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — Tech Stack
# ════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
navy_header(sl, "Technology Stack", "Production-grade tools across all layers")

layers = [
    ("Agent Layer",    NAVY,   [
        ("Band of Agents SDK (BandSDK)", "Multi-room P2P event bus for agent coordination"),
        ("LangGraph",                    "StateGraph workflows — deterministic agent execution"),
        ("Gemini 1.5 Flash / 2.0 Flash Lite", "LLM for PatientMgmtAgent & SummaryAgent"),
        ("Python asyncio",               "Non-blocking concurrency with Future-based proxy tools"),
    ]),
    ("Backend Layer",  TEAL,   [
        ("Express.js (TypeScript)",      "REST API server for frontend ↔ Supabase communication"),
        ("Supabase (PostgreSQL)",        "Patient records, appointments, prescriptions, inventory"),
        ("Row Level Security (RLS)",     "Per-user data isolation enforced at DB level"),
        ("Docker",                       "Containerized backend deployment"),
    ]),
    ("Frontend Layer", GREEN,  [
        ("Vite + TypeScript (React)",    "Doctor & Pharmacist Desktop dashboard (CarePulse)"),
        ("Mobile Interface (Web App)",   "Patient-facing chatbot & wayfinding app"),
        ("Socket.IO / WebSockets",       "Real-time queue & event streaming to UI"),
        ("Vercel / Cloudflare Tunnel",   "CDN deployment & secure tunnel for development"),
    ]),
]

lx = 0.3
for layer, color, items in layers:
    rect(sl, lx, 1.25, 4.15, 5.95, fill=PALE_GREY, line=color, line_w=Pt(2))
    rect(sl, lx, 1.25, 4.15, 0.48, fill=color)
    text_box(sl, layer, lx, 1.28, 4.15, 0.42,
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    iy = 1.83
    for tech, desc in items:
        rect(sl, lx + 0.15, iy, 3.85, 1.1, fill=WHITE, line=color, line_w=Pt(1))
        text_box(sl, tech, lx + 0.25, iy + 0.05, 3.65, 0.38,
                 font_size=12, bold=True, color=color)
        text_box(sl, desc, lx + 0.25, iy + 0.46, 3.65, 0.58,
                 font_size=10.5, color=MID_GREY, wrap=True)
        iy += 1.22
    lx += 4.35


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — Impact & Key Features
# ════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
navy_header(sl, "Impact & Key Features", "What M.A.S.H eliminates — and what it enables")

# Left: Eliminates
rect(sl, 0.3, 1.3, 6.0, 5.95, fill=RGBColor(0xFF,0xEB,0xEB), line=RED, line_w=Pt(2))
rect(sl, 0.3, 1.3, 6.0, 0.48, fill=RED)
text_box(sl, "What M.A.S.H Eliminates", 0.3, 1.33, 6.0, 0.42,
         font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

eliminates = [
    "Manual symptom-to-specialist matching",
    "Schedulers managing appointment books",
    "EHR tab-switching to build patient context",
    "Prescription written without stock check",
    "Patients wandering hallways for directions",
    "Pharmacist unaware of incoming prescriptions",
    "Last-minute drug unavailability surprises",
    "Hard-coded single-room agent bottlenecks",
]
ey = 1.9
for e in eliminates:
    text_box(sl, "✗  " + e, 0.5, ey, 5.6, 0.44, font_size=13, color=RED)
    ey += 0.55

# Right: Enables
rect(sl, 7.0, 1.3, 6.0, 5.95, fill=RGBColor(0xE8,0xF9,0xEE), line=GREEN, line_w=Pt(2))
rect(sl, 7.0, 1.3, 6.0, 0.48, fill=GREEN)
text_box(sl, "What M.A.S.H Enables", 7.0, 1.33, 6.0, 0.42,
         font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

enables = [
    "Instant AI-driven specialist matching",
    "Remote booking from anywhere, any time",
    "One-glance clinical summary on arrival",
    "Prescription safety audit before leaving room",
    "Step-by-step in-clinic navigation on phone",
    "Proactive pharmacy stock alerts",
    "Human-in-loop approval for out-of-stock drugs",
    "Horizontally scalable agent instances",
]
ney = 1.9
for e in enables:
    text_box(sl, "✓  " + e, 7.2, ney, 5.6, 0.44, font_size=13, color=GREEN)
    ney += 0.55

# Centre divider
text_box(sl, "VS", 6.3, 4.1, 0.73, 0.6, font_size=22, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 14 — Summary & What's Next
# ════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
rect(sl, 0, 6.8, 13.33, 0.7, fill=TEAL)

text_box(sl, "M.A.S.H", 0.5, 0.3, 12.33, 1.1,
         font_size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(sl, 4.0, 1.35, 5.33, 0.07, fill=TEAL)
text_box(sl, "Medical Assistant & Services Hub", 0.5, 1.45, 12.33, 0.55,
         font_size=22, color=LIGHT_TEAL, align=PP_ALIGN.CENTER)

# Summary bullets
summary = [
    ("6 AI Agents",       "Each owns one clinical domain — Registration, Navigation,\nSummary, Prescription, Stock, Patient Management"),
    ("6 Virtual Rooms",   "Persistent secure event buses — no dynamic creation overhead.\nOne concern per room, zero coupling between rooms"),
    ("Proxy Tool Bridge", "LLM agents + DB agents communicate via asyncio.Future +\nrequestId — decoupled, timeout-safe, horizontally scalable"),
    ("Safety First",      "LangGraph static graphs guarantee prescription safety gates\ncannot be bypassed. Human-in-loop built into the graph"),
]

sy4 = 2.15
for title, desc in summary:
    rect(sl, 0.6, sy4, 12.1, 0.98, fill=RGBColor(0x0D,0x2B,0x5A), line=TEAL, line_w=Pt(1.5))
    text_box(sl, title, 0.75, sy4 + 0.08, 2.8, 0.38,
             font_size=14, bold=True, color=TEAL)
    text_box(sl, desc, 3.65, sy4 + 0.08, 8.8, 0.82,
             font_size=12, color=WHITE, wrap=True)
    sy4 += 1.08

text_box(sl, "Built with Band of Agents SDK  ·  Supabase  ·  LangGraph  ·  Gemini  ·  Express.js",
         0.5, 6.82, 12.33, 0.4,
         font_size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────
out_path = "/home/bluni/dev/Band-of-agents/M.A.S.H/MASH_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}  ({len(prs.slides)} slides)")
